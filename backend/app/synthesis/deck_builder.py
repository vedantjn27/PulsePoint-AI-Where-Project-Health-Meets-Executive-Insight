"""Executive PowerPoint deck generation using python-pptx."""

from __future__ import annotations

from datetime import datetime
from pathlib import Path
import textwrap

from pptx import Presentation
from pptx.chart.data import CategoryChartData, ChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

from app.schemas.synthesis import DeckBrandingConfig, ProjectHealthSummary, SynthesisResponse
from app.synthesis import deck_theme as theme
from app.synthesis.consistency_check import validate_synthesis_consistency


OUTPUT_DIR = Path(__file__).resolve().parents[2] / "outputs" / "decks"
SLIDE_W = 13.333
SLIDE_H = 7.5
CONTENT_LEFT = 0.75
CONTENT_TOP = 1.25
CONTENT_W = 11.85


def build_monthly_deck(synthesis: SynthesisResponse, branding: DeckBrandingConfig | None = None) -> Path:
    validate_synthesis_consistency(synthesis)
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    brand = _resolve_branding(branding)
    suffix = _safe_suffix(brand["client_name"]) if brand["client_name"] else "default"
    run_stamp = datetime.now().strftime("%H%M%S")
    deck_path = OUTPUT_DIR / f"portfolio_health_review_{synthesis.generated_date.isoformat()}_{run_stamp}_{suffix}.pptx"

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    _title_slide(prs, synthesis, brand)
    _portfolio_snapshot_slide(prs, synthesis, brand)
    _trend_slide(prs, synthesis, brand)
    _risk_slide(prs, synthesis, brand)
    _intervention_map_slide(prs, synthesis, brand)
    _agent_insights_slide(prs, synthesis, brand)
    _action_plan_slide(prs, synthesis, brand)

    prs.save(deck_path)
    return deck_path


def _title_slide(prs: Presentation, synthesis: SynthesisResponse, brand: dict[str, str | None]) -> None:
    slide = _new_slide(prs, brand)
    client = brand["client_name"] or "Executive Portfolio"
    headline = f"{client} Health Review"
    verdict = _executive_verdict(synthesis)

    _cover_backdrop(slide, brand)
    _add_logo(slide, brand)
    _text(slide, "PulsePoint AI", 0.68, 0.36, 2.6, 0.28, 12, theme.WHITE, bold=True)
    _text(slide, "Monthly portfolio synthesis", 10.0, 0.36, 2.65, 0.22, 10, theme.WHITE, bold=True, align=PP_ALIGN.RIGHT)

    _text(slide, headline, 0.72, 1.62, 8.9, 0.92, 38, brand["primary_color"], bold=True)
    _text(slide, synthesis.period, 0.76, 2.48, 4.8, 0.35, 17, theme.GRAY)
    _cover_verdict(slide, verdict, 0.76, 3.08, 7.15, 1.45, brand)

    total_red = synthesis.rag_distribution.get("Red", 0)
    total_amber = synthesis.rag_distribution.get("Amber", 0)
    total_green = synthesis.rag_distribution.get("Green", 0)
    _cover_metric_card(slide, "Projects", str(synthesis.total_projects), "in portfolio", 0.76, 4.98, brand["accent_color"])
    _cover_metric_card(slide, "RAG Mix", f"{total_green}/{total_amber}/{total_red}", "green / amber / red", 2.72, 4.98, theme.RED if total_red else theme.GREEN)
    _cover_metric_card(slide, "Confidence", f"{synthesis.average_confidence:.0%}", "average data quality", 4.68, 4.98, brand["accent_color"])
    _cover_metric_card(slide, "Trend", synthesis.portfolio_trend.title(), "portfolio movement", 6.64, 4.98, _trend_color(synthesis.portfolio_trend))
    _cover_signal_strip(slide, synthesis, 8.58, 5.05, 3.72, 1.12, brand)
    _text(slide, "Deterministic scoring | Evidence-backed agent reasoning | Client-ready executive recommendations", 0.76, 6.62, 8.9, 0.28, 10, theme.GRAY)
    _footer(slide, synthesis, brand)


def _portfolio_snapshot_slide(prs: Presentation, synthesis: SynthesisResponse, brand: dict[str, str | None]) -> None:
    slide = _content_slide(prs, "Portfolio Health Snapshot", "Where the portfolio stands today", synthesis, brand)
    _rag_donut(slide, synthesis, 0.75, 1.35, 3.55, 2.85)
    _signal_bars(slide, synthesis, 4.65, 1.35, 3.35, 2.85)
    _project_scorecards(slide, synthesis.project_health[:4], 0.75, 4.62, 11.85, 1.65, brand)
    _insight_strip(
        slide,
        [
            f"{synthesis.rag_distribution.get('Red', 0)} project(s) need executive attention.",
            f"{len([p for p in synthesis.project_health if p.data_confidence < 0.7])} project(s) have confidence below 70%.",
            _strongest_theme(synthesis),
        ],
        8.42,
        1.35,
        3.85,
        brand,
    )


def _trend_slide(prs: Presentation, synthesis: SynthesisResponse, brand: dict[str, str | None]) -> None:
    slide = _content_slide(prs, "Cross-Project Trends", "Movement over time, not isolated project summaries", synthesis, brand)
    _line_chart(slide, synthesis, 0.75, 1.35, 7.1, 3.3)
    _rag_stacked_chart(slide, synthesis, 8.25, 1.35, 4.0, 3.3)
    trend_note = {
        "improving": "Average health is improving, but Red concentration still requires active management.",
        "declining": "Average health is declining; leadership should intervene before risk becomes contractual.",
        "stable": "Average health is stable; focus should shift to weak signals and outlier projects.",
    }.get(synthesis.portfolio_trend, "Trend is stable based on available snapshots.")
    _callout(slide, "Trend interpretation", trend_note, 0.75, 5.1, 5.75, 1.05, _trend_color(synthesis.portfolio_trend))
    _callout(slide, "What changed", _movement_sentence(synthesis), 6.8, 5.1, 5.45, 1.05, brand["accent_color"])


def _risk_slide(prs: Presentation, synthesis: SynthesisResponse, brand: dict[str, str | None]) -> None:
    slide = _content_slide(prs, "Emerging Risk Themes", "Patterns requiring leadership action", synthesis, brand)
    themes = synthesis.systemic_themes or ["No systemic delivery theme is currently visible."]
    risks = synthesis.emerging_risks or ["No emerging risk has been detected from the latest scored snapshots."]
    for idx, item in enumerate(themes[:2]):
        _risk_card(slide, f"Theme {idx + 1}", item, 0.75 + idx * 5.9, 1.35, 5.45, 1.15, _theme_color(item))
    _risk_matrix(slide, synthesis, 0.75, 3.0, 5.8, 3.15, brand)
    _ranked_list(slide, "Top emerging risks", risks[:4], 6.9, 3.0, 5.35, 3.15, brand["accent_color"], max_items=4)


def _intervention_map_slide(prs: Presentation, synthesis: SynthesisResponse, brand: dict[str, str | None]) -> None:
    slide = _content_slide(prs, "Executive Intervention Map", "Prioritize action by health score and evidence confidence", synthesis, brand)
    _scatter_map(slide, synthesis.project_health, 0.75, 1.25, 6.45, 4.9, brand)
    critical = [p for p in synthesis.project_health if p.rag_status == "Red"][:3]
    watch = [p for p in synthesis.project_health if p.rag_status == "Amber"][:2]
    healthy = [p for p in synthesis.project_health if p.rag_status == "Green"][:2]
    _ranked_list(slide, "Intervene now", [_project_action_line(p) for p in critical] or ["No Red projects in current portfolio."], 7.55, 1.25, 4.85, 1.55, theme.RED, max_items=2)
    _ranked_list(slide, "Watch closely", [_project_action_line(p) for p in watch] or ["No Amber projects in current portfolio."], 7.55, 3.05, 4.85, 1.25, theme.AMBER, max_items=2)
    _ranked_list(slide, "Protect momentum", [_project_action_line(p) for p in healthy] or ["No Green projects in current portfolio."], 7.55, 4.6, 4.85, 1.25, theme.GREEN, max_items=2)


def _agent_insights_slide(prs: Presentation, synthesis: SynthesisResponse, brand: dict[str, str | None]) -> None:
    slide = _content_slide(prs, "Agent Reasoning And Evidence", "What the scoring agent surfaced across plans", synthesis, brand)
    risk_lines = _agent_risk_lines(synthesis)
    action_lines = _agent_action_lines(synthesis)
    _evidence_panel(slide, "Evidence-backed risk signals", risk_lines[:3], 0.75, 1.35, 5.75, 4.55, theme.RED)
    _evidence_panel(slide, "Recommended management actions", action_lines[:3], 6.85, 1.35, 5.5, 4.55, brand["accent_color"])
    _callout(
        slide,
        "Governance note",
        "The agent explains why scores moved and records the reasoning trace; deterministic scoring remains the source of truth.",
        0.75,
        6.35,
        11.6,
        0.65,
        brand["primary_color"],
    )


def _action_plan_slide(prs: Presentation, synthesis: SynthesisResponse, brand: dict[str, str | None]) -> None:
    slide = _content_slide(prs, "30-Day Executive Action Plan", "Client-ready decisions and next steps", synthesis, brand)
    recommendations = synthesis.recommendations or ["Maintain weekly governance and continue monitoring project health."]
    columns = [
        ("This week", recommendations[0:2] or ["Confirm portfolio owners and escalation paths."]),
        ("Next 2 weeks", recommendations[2:4] or ["Run focused reviews on weak signals and open blockers."]),
        ("By month end", ["Refresh project plans, regenerate synthesis, and confirm RAG movement with evidence."]),
    ]
    for idx, (title, items) in enumerate(columns):
        _execution_stage(slide, title, items, 0.75 + idx * 4.05, 1.35, 3.65, 4.45, brand, idx)
    _methodology_band(slide, 0.75, 6.12, 11.7, 0.62, brand)


def _new_slide(prs: Presentation, brand: dict[str, str | None]):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background.fill
    background.solid()
    background.fore_color.rgb = RGBColor.from_string(theme.WHITE)
    _side_rule(slide, brand)
    return slide


def _content_slide(prs: Presentation, title: str, subtitle: str, synthesis: SynthesisResponse, brand: dict[str, str | None]):
    slide = _new_slide(prs, brand)
    _text(slide, title, 0.75, 0.35, 7.6, 0.42, 25, brand["primary_color"], bold=True)
    _text(slide, subtitle, 0.77, 0.83, 7.8, 0.25, 10, theme.GRAY)
    _footer(slide, synthesis, brand)
    return slide


def _side_rule(slide, brand: dict[str, str | None]) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.12), Inches(SLIDE_H))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor.from_string(brand["accent_color"])
    shape.line.fill.background()


def _accent_block(slide, brand: dict[str, str | None], x: float, y: float, w: float, h: float) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor.from_string(brand["primary_color"])
    shape.line.fill.background()


def _cover_backdrop(slide, brand: dict[str, str | None]) -> None:
    _accent_block(slide, brand, 0.0, 0.0, SLIDE_W, 1.05)
    hero = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.58), Inches(1.55), Inches(3.72), Inches(3.15))
    hero.fill.solid()
    hero.fill.fore_color.rgb = RGBColor.from_string("F8FAFD")
    hero.line.color.rgb = RGBColor.from_string("D8DEE8")
    _text(slide, "Portfolio command view", 8.9, 1.88, 2.8, 0.24, 12, brand["primary_color"], bold=True)
    _text(slide, "Score movement, RAG pressure, and executive actions in one client-ready review.", 8.9, 2.25, 2.85, 0.58, 10, theme.GRAY)
    for idx, (label, color, width) in enumerate(
        [
            ("Health trend", theme.GREEN, 1.85),
            ("Risk pressure", theme.RED, 1.25),
            ("Action clarity", brand["accent_color"], 1.55),
        ]
    ):
        yy = 3.05 + idx * 0.42
        _text(slide, label, 8.9, yy - 0.02, 1.1, 0.16, 8, theme.NAVY, bold=True)
        rail = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(10.15), Inches(yy), Inches(1.9), Inches(0.14))
        rail.fill.solid()
        rail.fill.fore_color.rgb = RGBColor.from_string("E2E8F0")
        rail.line.fill.background()
        bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(10.15), Inches(yy), Inches(width), Inches(0.14))
        bar.fill.solid()
        bar.fill.fore_color.rgb = RGBColor.from_string(color)
        bar.line.fill.background()


def _cover_verdict(slide, verdict: str, x: float, y: float, w: float, h: float, brand: dict[str, str | None]) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor.from_string(theme.LIGHT_GRAY)
    shape.line.color.rgb = RGBColor.from_string(brand["accent_color"])
    shape.line.width = Pt(1.5)
    _text(slide, "Executive readout", x + 0.28, y + 0.18, w - 0.56, 0.28, 13, brand["accent_color"], bold=True)
    _text(slide, _wrap(verdict, _wrap_width(w - 0.56), max_lines=3), x + 0.28, y + 0.54, w - 0.56, 0.62, 15, theme.NAVY)


def _cover_metric_card(slide, label: str, value: str, hint: str, x: float, y: float, color: str) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(1.66), Inches(1.18))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor.from_string(theme.WHITE)
    shape.line.color.rgb = RGBColor.from_string("D8DEE8")
    _text(slide, _clip(value, 12), x + 0.16, y + 0.16, 1.34, 0.34, 20, color, bold=True)
    _text(slide, label, x + 0.16, y + 0.56, 1.25, 0.22, 10, theme.NAVY, bold=True)
    _text(slide, hint, x + 0.16, y + 0.82, 1.38, 0.24, 8, theme.GRAY)


def _cover_signal_strip(slide, synthesis: SynthesisResponse, x: float, y: float, w: float, h: float, brand: dict[str, str | None]) -> None:
    strip = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    strip.fill.solid()
    strip.fill.fore_color.rgb = RGBColor.from_string("F8FAFD")
    strip.line.color.rgb = RGBColor.from_string("D8DEE8")
    strongest = _strongest_theme(synthesis)
    items = [
        ("Top theme", _clip(strongest, 56), brand["accent_color"]),
        ("Priority", _clip((synthesis.recommendations or ["Continue weekly governance."])[0], 60), theme.RED if synthesis.rag_distribution.get("Red", 0) else theme.GREEN),
        ("Data quality", f"{synthesis.average_confidence:.0%} average confidence", brand["primary_color"]),
    ]
    for idx, (label, value, color) in enumerate(items):
        yy = y + 0.14 + idx * 0.32
        _text(slide, label, x + 0.22, yy, 0.9, 0.16, 8, theme.GRAY, bold=True)
        _text(slide, _clip(value, 38), x + 1.1, yy - 0.01, w - 1.3, 0.17, 8, color, bold=True)


def _text(slide, text: str, x: float, y: float, w: float, h: float, size: int, color: str, *, bold: bool = False, align=PP_ALIGN.LEFT) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.02)
    frame.margin_right = Inches(0.02)
    frame.margin_top = Inches(0.01)
    frame.margin_bottom = Inches(0.01)
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.name = theme.FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor.from_string(color)


def _callout(slide, title: str, body: str, x: float, y: float, w: float, h: float, color: str) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor.from_string(theme.LIGHT_GRAY)
    shape.line.color.rgb = RGBColor.from_string(color)
    shape.line.width = Pt(1.2)
    _text(slide, title, x + 0.22, y + 0.14, w - 0.44, 0.25, 11, color, bold=True)
    _text(slide, _wrap(body, _wrap_width(w), max_lines=3), x + 0.22, y + 0.46, w - 0.44, h - 0.58, 11, theme.NAVY)


def _metric_card(slide, label: str, value: str, hint: str, x: float, y: float, color: str) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(1.7), Inches(1.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor.from_string(theme.WHITE)
    shape.line.color.rgb = RGBColor.from_string("D8DEE8")
    _text(slide, _clip(value, 16), x + 0.16, y + 0.18, 1.35, 0.34, 20, color, bold=True)
    _text(slide, label, x + 0.16, y + 0.58, 1.3, 0.22, 9, theme.NAVY, bold=True)
    _text(slide, hint, x + 0.16, y + 0.83, 1.38, 0.24, 8, theme.GRAY)


def _rag_donut(slide, synthesis: SynthesisResponse, x: float, y: float, w: float, h: float) -> None:
    data = ChartData()
    data.categories = ["Green", "Amber", "Red"]
    data.add_series("Projects", [synthesis.rag_distribution.get("Green", 0), synthesis.rag_distribution.get("Amber", 0), synthesis.rag_distribution.get("Red", 0)])
    chart = slide.shapes.add_chart(XL_CHART_TYPE.DOUGHNUT, Inches(x), Inches(y), Inches(w), Inches(h), data).chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.plots[0].has_data_labels = True
    _text(slide, "RAG distribution", x, y - 0.28, w, 0.2, 10, theme.NAVY, bold=True)


def _signal_bars(slide, synthesis: SynthesisResponse, x: float, y: float, w: float, h: float) -> None:
    data = CategoryChartData()
    data.categories = [item.signal for item in synthesis.signal_health]
    data.add_series("Average score", [item.average_score or 0 for item in synthesis.signal_health])
    chart = slide.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, Inches(x), Inches(y), Inches(w), Inches(h), data).chart
    chart.has_legend = False
    chart.value_axis.maximum_scale = 100
    chart.value_axis.minimum_scale = 0
    chart.category_axis.tick_labels.font.size = Pt(8)
    chart.value_axis.tick_labels.font.size = Pt(8)
    _text(slide, "Average signal health", x, y - 0.28, w, 0.2, 10, theme.NAVY, bold=True)


def _project_table(slide, projects: list[ProjectHealthSummary], x: float, y: float, w: float, h: float) -> None:
    _text(slide, "Project scoreboard", x, y - 0.28, w, 0.2, 10, theme.NAVY, bold=True)
    rows = max(2, len(projects) + 1)
    table = slide.shapes.add_table(rows, 5, Inches(x), Inches(y), Inches(w), Inches(h)).table
    widths = [3.6, 1.0, 1.15, 1.2, 4.9]
    for idx, width in enumerate(widths):
        table.columns[idx].width = Inches(width)
    headers = ["Project", "RAG", "Score", "Confidence", "Primary risk / action"]
    for col, header in enumerate(headers):
        _cell(table.cell(0, col), header, theme.WHITE, theme.NAVY, bold=True, size=8)
    for row_idx, project in enumerate(projects, start=1):
        risk = project.top_risks[0] if project.top_risks else _project_action_line(project)
        values = [project.project_name, project.rag_status, f"{project.composite_score:.1f}", f"{project.data_confidence:.0%}", risk]
        for col, value in enumerate(values):
            fill = _rag_color(project.rag_status) if col == 1 else theme.WHITE
            color = theme.WHITE if col == 1 else theme.NAVY
            _cell(table.cell(row_idx, col), _wrap(str(value), 58), color, fill, bold=col == 1, size=8)


def _project_scorecards(slide, projects: list[ProjectHealthSummary], x: float, y: float, w: float, h: float, brand: dict[str, str | None]) -> None:
    _text(slide, "Project scoreboard", x, y - 0.28, w, 0.22, 11, theme.NAVY, bold=True)
    if not projects:
        _callout(slide, "No project snapshots", "Run portfolio analysis to populate project-level health cards.", x, y, w, h, brand["accent_color"])
        return
    gap = 0.18
    card_w = (w - gap * (len(projects) - 1)) / len(projects)
    for idx, project in enumerate(projects):
        cx = x + idx * (card_w + gap)
        color = _rag_color(project.rag_status)
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(cx), Inches(y), Inches(card_w), Inches(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor.from_string(theme.WHITE)
        shape.line.color.rgb = RGBColor.from_string("D8DEE8")
        _status_pill(slide, project.rag_status, cx + 0.16, y + 0.14, color)
        _text(slide, _clip(project.project_name, 24), cx + 0.16, y + 0.48, card_w - 0.32, 0.25, 10, theme.NAVY, bold=True)
        _text(slide, f"{project.composite_score:.1f} score | {project.data_confidence:.0%} confidence", cx + 0.16, y + 0.8, card_w - 0.32, 0.22, 8, theme.GRAY)
        risk = project.top_risks[0] if project.top_risks else "No material risk surfaced."
        _text(slide, _wrap(risk, _wrap_width(card_w - 0.32), max_lines=2), cx + 0.16, y + 1.08, card_w - 0.32, 0.42, 8, brand["primary_color"])


def _status_pill(slide, label: str, x: float, y: float, color: str) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(0.78), Inches(0.25))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor.from_string(color)
    shape.line.fill.background()
    _text(slide, label, x + 0.04, y + 0.04, 0.7, 0.14, 8, theme.WHITE, bold=True, align=PP_ALIGN.CENTER)


def _insight_strip(slide, items: list[str], x: float, y: float, w: float, brand: dict[str, str | None]) -> None:
    _text(slide, "Executive implications", x, y - 0.28, w, 0.2, 10, theme.NAVY, bold=True)
    for idx, item in enumerate(items[:3]):
        _callout(slide, f"Insight {idx + 1}", item, x, y + idx * 1.03, w, 0.82, brand["accent_color"])


def _line_chart(slide, synthesis: SynthesisResponse, x: float, y: float, w: float, h: float) -> None:
    data = CategoryChartData()
    data.categories = [point.run_date.strftime("%d %b") for point in synthesis.trend_points]
    data.add_series("Average score", [point.average_score for point in synthesis.trend_points])
    chart = slide.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS, Inches(x), Inches(y), Inches(w), Inches(h), data).chart
    chart.has_legend = False
    chart.value_axis.maximum_scale = 100
    chart.value_axis.minimum_scale = 0
    chart.value_axis.tick_labels.font.size = Pt(8)
    chart.category_axis.tick_labels.font.size = Pt(8)
    _text(slide, "Average portfolio score", x, y - 0.28, w, 0.2, 10, theme.NAVY, bold=True)


def _rag_stacked_chart(slide, synthesis: SynthesisResponse, x: float, y: float, w: float, h: float) -> None:
    data = CategoryChartData()
    data.categories = [point.run_date.strftime("%d %b") for point in synthesis.trend_points]
    data.add_series("Green", [point.rag_counts.get("Green", 0) for point in synthesis.trend_points])
    data.add_series("Amber", [point.rag_counts.get("Amber", 0) for point in synthesis.trend_points])
    data.add_series("Red", [point.rag_counts.get("Red", 0) for point in synthesis.trend_points])
    chart = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_STACKED, Inches(x), Inches(y), Inches(w), Inches(h), data).chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.value_axis.tick_labels.font.size = Pt(8)
    chart.category_axis.tick_labels.font.size = Pt(8)
    _text(slide, "RAG movement", x, y - 0.28, w, 0.2, 10, theme.NAVY, bold=True)


def _risk_card(slide, label: str, text: str, x: float, y: float, w: float, h: float, color: str) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor.from_string(theme.WHITE)
    shape.line.color.rgb = RGBColor.from_string(color)
    shape.line.width = Pt(1.1)
    _text(slide, label, x + 0.18, y + 0.15, w - 0.36, 0.22, 10, color, bold=True)
    _text(slide, _wrap(text, _wrap_width(w - 0.36), max_lines=3), x + 0.18, y + 0.45, w - 0.36, h - 0.55, 10, theme.NAVY)


def _evidence_panel(slide, title: str, items: list[str], x: float, y: float, w: float, h: float, color: str) -> None:
    panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    panel.fill.solid()
    panel.fill.fore_color.rgb = RGBColor.from_string(theme.WHITE)
    panel.line.color.rgb = RGBColor.from_string("D8DEE8")
    panel.line.width = Pt(1.1)
    _text(slide, title, x + 0.28, y + 0.22, w - 0.56, 0.32, 14, color, bold=True)

    visible = items[:3] or ["No material signal was available."]
    row_h = 1.04
    for idx, item in enumerate(visible):
        yy = y + 0.82 + idx * 1.16
        row = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x + 0.28), Inches(yy), Inches(w - 0.56), Inches(row_h))
        row.fill.solid()
        row.fill.fore_color.rgb = RGBColor.from_string("F8FAFD")
        row.line.color.rgb = RGBColor.from_string("E2E8F0")
        marker = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 0.48), Inches(yy + 0.28), Inches(0.34), Inches(0.34))
        marker.fill.solid()
        marker.fill.fore_color.rgb = RGBColor.from_string(color)
        marker.line.fill.background()
        _text(slide, str(idx + 1), x + 0.48, yy + 0.365, 0.34, 0.13, 9, theme.WHITE, bold=True, align=PP_ALIGN.CENTER)
        _text(slide, _wrap(item, _wrap_width(w - 1.35), max_lines=2), x + 0.98, yy + 0.24, w - 1.5, 0.55, 11, theme.NAVY)


def _risk_matrix(slide, synthesis: SynthesisResponse, x: float, y: float, w: float, h: float, brand: dict[str, str | None]) -> None:
    _text(slide, "Risk concentration matrix", x, y - 0.28, w, 0.2, 10, theme.NAVY, bold=True)
    base = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    base.fill.solid()
    base.fill.fore_color.rgb = RGBColor.from_string("F8FAFD")
    base.line.color.rgb = RGBColor.from_string("D8DEE8")
    labels = [("Low urgency", x + 0.2, y + h - 0.38), ("High urgency", x + w - 1.3, y + 0.13), ("Low evidence", x + 0.2, y + 0.13), ("High evidence", x + w - 1.35, y + h - 0.38)]
    for label, lx, ly in labels:
        _text(slide, label, lx, ly, 1.2, 0.18, 8, theme.GRAY)
    points = synthesis.project_health[:8]
    for project in points:
        px = x + 0.45 + (project.data_confidence * (w - 1.0))
        urgency = max(0, min(1, (100 - project.composite_score) / 100))
        py = y + 0.35 + (urgency * (h - 0.85))
        bubble = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(px), Inches(py), Inches(0.28), Inches(0.28))
        bubble.fill.solid()
        bubble.fill.fore_color.rgb = RGBColor.from_string(_rag_color(project.rag_status))
        bubble.line.color.rgb = RGBColor.from_string(theme.WHITE)
    _text(slide, "Each dot is a project: right = stronger evidence, higher = greater intervention need.", x + 0.25, y + h - 0.78, w - 0.5, 0.18, 8, brand["primary_color"])


def _scatter_map(slide, projects: list[ProjectHealthSummary], x: float, y: float, w: float, h: float, brand: dict[str, str | None]) -> None:
    _text(slide, "Score vs. confidence", x, y - 0.28, w, 0.2, 10, theme.NAVY, bold=True)
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor.from_string("F8FAFD")
    box.line.color.rgb = RGBColor.from_string("D8DEE8")
    _axis_label(slide, "Lower score", x + 0.15, y + h - 0.35)
    _axis_label(slide, "Higher score", x + 0.15, y + 0.15)
    _axis_label(slide, "Higher confidence", x + w - 1.45, y + h - 0.35)
    for project in projects[:10]:
        px = x + 0.45 + (project.data_confidence * (w - 1.0))
        py = y + 0.35 + ((100 - project.composite_score) / 100 * (h - 0.8))
        bubble = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(px), Inches(py), Inches(0.36), Inches(0.36))
        bubble.fill.solid()
        bubble.fill.fore_color.rgb = RGBColor.from_string(_rag_color(project.rag_status))
        bubble.line.color.rgb = RGBColor.from_string(theme.WHITE)
        _text(slide, _short_name(project.project_name), px + 0.42, py + 0.0, 1.55, 0.2, 8, brand["primary_color"])


def _ranked_list(slide, title: str, items: list[str], x: float, y: float, w: float, h: float, color: str, *, max_items: int = 4) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor.from_string(theme.WHITE)
    shape.line.color.rgb = RGBColor.from_string("D8DEE8")
    _text(slide, title, x + 0.22, y + 0.18, w - 0.44, 0.28, 12, color, bold=True)
    visible = items[:max_items] or ["No item available."]
    row_h = (h - 0.72) / max(1, len(visible))
    for idx, item in enumerate(visible):
        yy = y + 0.58 + idx * row_h
        marker = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 0.22), Inches(yy + 0.02), Inches(0.25), Inches(0.25))
        marker.fill.solid()
        marker.fill.fore_color.rgb = RGBColor.from_string(color)
        marker.line.fill.background()
        _text(slide, str(idx + 1), x + 0.22, yy + 0.055, 0.25, 0.13, 8, theme.WHITE, bold=True, align=PP_ALIGN.CENTER)
        _text(slide, _wrap(item, _wrap_width(w - 0.8), max_lines=2), x + 0.55, yy, w - 0.8, min(0.48, row_h - 0.04), 9, theme.NAVY)


def _plan_column(slide, title: str, items: list[str], x: float, y: float, w: float, h: float, brand: dict[str, str | None], index: int) -> None:
    colors = [brand["accent_color"], theme.AMBER, theme.GREEN]
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor.from_string(theme.WHITE)
    shape.line.color.rgb = RGBColor.from_string(colors[index])
    _text(slide, title, x + 0.22, y + 0.22, w - 0.44, 0.3, 14, colors[index], bold=True)
    for idx, item in enumerate(items[:3]):
        yy = y + 0.83 + idx * 0.9
        _text(slide, f"{idx + 1}", x + 0.25, yy, 0.26, 0.2, 10, colors[index], bold=True)
        _text(slide, _wrap(item, _wrap_width(w - 0.9), max_lines=2), x + 0.62, yy - 0.02, w - 0.9, 0.58, 10, theme.NAVY)


def _execution_stage(slide, title: str, items: list[str], x: float, y: float, w: float, h: float, brand: dict[str, str | None], index: int) -> None:
    colors = [brand["accent_color"], theme.AMBER, theme.GREEN]
    color = colors[index]
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor.from_string(theme.WHITE)
    shape.line.color.rgb = RGBColor.from_string("D8DEE8")
    shape.line.width = Pt(1.0)

    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.18))
    band.fill.solid()
    band.fill.fore_color.rgb = RGBColor.from_string(color)
    band.line.fill.background()

    phase = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 0.25), Inches(y + 0.42), Inches(0.48), Inches(0.48))
    phase.fill.solid()
    phase.fill.fore_color.rgb = RGBColor.from_string(color)
    phase.line.fill.background()
    _text(slide, str(index + 1), x + 0.25, y + 0.56, 0.48, 0.14, 10, theme.WHITE, bold=True, align=PP_ALIGN.CENTER)
    _text(slide, title, x + 0.88, y + 0.42, w - 1.1, 0.34, 16, color, bold=True)
    _text(slide, _stage_subtitle(index), x + 0.88, y + 0.82, w - 1.1, 0.25, 9, theme.GRAY)

    visible = items[:2] or ["Confirm next governance action."]
    for item_idx, item in enumerate(visible):
        yy = y + 1.42 + item_idx * 1.28
        action = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x + 0.28), Inches(yy), Inches(w - 0.56), Inches(1.0))
        action.fill.solid()
        action.fill.fore_color.rgb = RGBColor.from_string("F8FAFD")
        action.line.color.rgb = RGBColor.from_string("E2E8F0")
        _text(slide, f"Action {item_idx + 1}", x + 0.52, yy + 0.17, 0.82, 0.18, 8, color, bold=True)
        _text(slide, _wrap(item, _wrap_width(w - 1.0), max_lines=2), x + 0.52, yy + 0.43, w - 1.0, 0.42, 11, theme.NAVY)

    _text(slide, _stage_outcome(index), x + 0.35, y + h - 0.56, w - 0.7, 0.28, 10, brand["primary_color"], bold=True)


def _methodology_band(slide, x: float, y: float, w: float, h: float, brand: dict[str, str | None]) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor.from_string(theme.LIGHT_GRAY)
    shape.line.fill.background()
    text = "Scoring is deterministic and configurable. The agent explains evidence and recommendations; it does not override the score. Every deck generation is audit logged."
    _text(slide, _wrap(text, 150, max_lines=2), x + 0.25, y + 0.22, w - 0.5, h - 0.3, 11, brand["primary_color"])


def _cell(cell, text: str, font_color: str, fill_color: str, *, bold: bool = False, size: int = 8) -> None:
    cell.fill.solid()
    cell.fill.fore_color.rgb = RGBColor.from_string(fill_color)
    cell.margin_left = Inches(0.05)
    cell.margin_right = Inches(0.05)
    cell.margin_top = Inches(0.03)
    cell.margin_bottom = Inches(0.03)
    frame = cell.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = text
    run.font.name = theme.FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor.from_string(font_color)


def _footer(slide, synthesis: SynthesisResponse, brand: dict[str, str | None]) -> None:
    label = brand["client_name"] or "Portfolio Health Review"
    _text(slide, f"{label} | {synthesis.period}", 0.75, 7.08, 5.7, 0.16, 8, theme.GRAY)


def _add_logo(slide, brand: dict[str, str | None]) -> None:
    logo_path = brand.get("logo_path")
    if not logo_path:
        return
    path = Path(logo_path)
    if path.exists() and path.is_file():
        slide.shapes.add_picture(str(path), Inches(11.1), Inches(0.25), width=Inches(1.45))


def _resolve_branding(branding: DeckBrandingConfig | None) -> dict[str, str | None]:
    if branding is None or branding.use_default_branding:
        return {
            "client_name": None,
            "primary_color": theme.NAVY,
            "accent_color": theme.BLUE,
            "logo_path": None,
        }
    return {
        "client_name": branding.client_name,
        "primary_color": _strip_hash(branding.primary_color or theme.NAVY),
        "accent_color": _strip_hash(branding.accent_color or theme.BLUE),
        "logo_path": branding.logo_path,
    }


def _executive_verdict(synthesis: SynthesisResponse) -> str:
    red = synthesis.rag_distribution.get("Red", 0)
    amber = synthesis.rag_distribution.get("Amber", 0)
    if red:
        return f"The portfolio is moving {synthesis.portfolio_trend}, but {red} Red project(s) require named executive intervention before the next client checkpoint."
    if amber:
        return f"The portfolio is broadly controlled with {amber} Amber project(s) requiring closer weekly governance."
    return "The portfolio is healthy; leadership should protect momentum and continue evidence-based monitoring."


def _strongest_theme(synthesis: SynthesisResponse) -> str:
    return synthesis.systemic_themes[0] if synthesis.systemic_themes else "No dominant systemic risk theme is visible."


def _movement_sentence(synthesis: SynthesisResponse) -> str:
    points = synthesis.trend_points
    if len(points) < 2:
        return "Not enough trend history is available yet to compare movement."
    delta = points[-1].average_score - points[0].average_score
    return f"Average portfolio score moved {delta:+.1f} points from {points[0].run_date.strftime('%d %b')} to {points[-1].run_date.strftime('%d %b')}."


def _agent_risk_lines(synthesis: SynthesisResponse) -> list[str]:
    lines = []
    for project in synthesis.project_health:
        for risk in project.top_risks[:2]:
            lines.append(f"{project.project_name}: {risk}")
    return lines or synthesis.emerging_risks or ["No material agent risk signal was available."]


def _agent_action_lines(synthesis: SynthesisResponse) -> list[str]:
    lines = []
    for project in synthesis.project_health:
        for action in project.recommended_actions[:1]:
            lines.append(f"{project.project_name}: {action}")
    return lines or synthesis.recommendations


def _project_action_line(project: ProjectHealthSummary) -> str:
    action = project.recommended_actions[0] if project.recommended_actions else "Confirm next governance action."
    return f"{project.project_name} ({project.composite_score:.1f}): {action}"


def _theme_color(text: str) -> str:
    lowered = text.lower()
    if "blocker" in lowered or "critical" in lowered:
        return theme.RED
    if "budget" in lowered:
        return theme.AMBER
    if "schedule" in lowered or "milestone" in lowered:
        return theme.BLUE
    return theme.GRAY


def _trend_color(value: str) -> str:
    if value == "improving":
        return theme.GREEN
    if value == "declining":
        return theme.RED
    return theme.AMBER


def _stage_subtitle(index: int) -> str:
    return [
        "Stabilize immediate delivery risk",
        "Reset commitments and owners",
        "Validate movement with evidence",
    ][index]


def _stage_outcome(index: int) -> str:
    return [
        "Outcome: named interventions before the next checkpoint",
        "Outcome: visible recovery path and accountable owners",
        "Outcome: refreshed RAG view ready for client governance",
    ][index]


def _rag_color(status: str) -> str:
    return {"Green": theme.GREEN, "Amber": theme.AMBER, "Red": theme.RED}.get(status, theme.GRAY)


def _axis_label(slide, text: str, x: float, y: float) -> None:
    _text(slide, text, x, y, 1.3, 0.18, 8, theme.GRAY)


def _short_name(value: str) -> str:
    words = value.replace("-", " ").split()
    return " ".join(words[:2]) if words else value[:14]


def _wrap(value: str, width: int, *, max_lines: int = 3) -> str:
    cleaned = " ".join(str(value).split())
    if not cleaned:
        return ""
    return "\n".join(textwrap.wrap(cleaned, width=max(12, width), max_lines=max_lines, placeholder="..."))


def _wrap_width(box_width_inches: float) -> int:
    return max(22, int(box_width_inches * 13.5))


def _clip(value: str, limit: int) -> str:
    cleaned = " ".join(str(value).split())
    if len(cleaned) <= limit:
        return cleaned
    return cleaned[: max(0, limit - 3)].rstrip() + "..."


def _strip_hash(value: str) -> str:
    return value.strip().lstrip("#")


def _safe_suffix(value: str) -> str:
    return "".join(character.lower() if character.isalnum() else "_" for character in value).strip("_")[:40] or "custom"
